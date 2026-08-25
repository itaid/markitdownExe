"""转换核心单元测试（无 GUI 依赖，跨平台 CI 可跑）"""
from pathlib import Path

import pytest


def _make_test_image(path: Path, label: str):
    """生成一张 > 4KB、含可 OCR 文字的测试图（条纹提供熵以撑大小）"""
    from PIL import Image, ImageDraw
    im = Image.new("RGB", (820, 400), "white")
    d = ImageDraw.Draw(im)
    # 浅灰细条纹 → 提高 PNG 熵，撑过 4KB 阈值，且不影响红框内文字识别
    for y in range(0, 400, 6):
        d.line([(0, y), (240, y)], fill=(200, 200, 200), width=1)
    d.rectangle([0, 0, 819, 46], fill="black")
    d.text((16, 56), "Quarter report", fill="black")
    d.text((16, 96), f"Label {label}", fill="black")
    d.text((16, 136), "Q1=100 Q2=200 Q3=300", fill="black")
    d.text((16, 176), "Total 600 units", fill="black")
    im.save(path, optimize=False)


from doc2md_tool.service import (
    ConvertOptions,
    collect_files,
    convert_file,
    normalize_filename,
    split_by_headings,
)


MD_DOC = """# 员工手册

## 总则
本文档规定公司员工的基本行为准则。

## 考勤制度
- 工作时间 9:00-18:00
- 打卡方式：钉钉

### 请假流程
主管审批 -> HR备案

## 保密条款
员工不得泄露商业机密。
"""


def test_normalize_filename():
    assert normalize_filename("新建文件夹 (2) 最终版") == "新建文件夹"
    assert normalize_filename("a/b:c*d?") == "a_b_c_d"
    assert normalize_filename("   copy  ") == "copy"
    assert normalize_filename("   ") == "unnamed"
    assert len(normalize_filename("x" * 200)) <= 80


def test_split_by_headings():
    sections = split_by_headings(MD_DOC, 2)
    # 首标题 “# 员工手册” 之前无内容 → 不产生 preamble；level 1 标题也作为切分点
    titles = [t for t, _ in sections]
    assert titles == ["员工手册", "总则", "考勤制度", "保密条款"]
    # 每个小节带上自己的标题行
    assert sections[2][1].startswith("## 考勤制度")
    assert "请假流程" in sections[2][1]
    # level=1 时不按 ## 切
    assert len(split_by_headings(MD_DOC, 1)) == 1  # 全部归入 “员工手册”
    # level=0 不切
    assert len(split_by_headings(MD_DOC, 0)) == 1
    # 前面有引言内容时产生 preamble
    with_preamble = "引言段落\\n\\n" + MD_DOC
    assert split_by_headings(with_preamble, 2)[0][0] == "__preamble__"


def test_convert_md_file(tmp_path: Path):
    src = tmp_path / "测试文档.md"
    src.write_text(MD_DOC, encoding="utf-8")
    out = tmp_path / "output"
    r = convert_file(src, ConvertOptions(output_dir=out))
    assert r.status == "ok", r.message
    assert r.dst == [out / "测试文档.md"]
    content = r.dst[0].read_text(encoding="utf-8")
    assert "考勤制度" in content


def test_convert_split(tmp_path: Path):
    src = tmp_path / "制度.md"
    src.write_text(MD_DOC, encoding="utf-8")
    out = tmp_path / "output"
    r = convert_file(src, ConvertOptions(split_at_heading_level=2, output_dir=out))
    assert r.status == "ok", r.message
    assert len(r.dst) == 4  # preamble + 3 个二级标题
    assert (out / "制度").is_dir()
    for p in r.dst:
        assert p.exists()
        assert p.read_text(encoding="utf-8").strip()


def test_convert_report(tmp_path: Path):
    import csv

    src = tmp_path / "a.md"
    src.write_text("hello world test content", encoding="utf-8")
    out = tmp_path / "output"
    r = convert_file(src, ConvertOptions(output_dir=out))
    from doc2md_tool.service import write_report

    write_report([r], out / "report.csv")
    with (out / "report.csv").open(encoding="utf-8-sig") as fh:
        rows = list(csv.reader(fh))
    assert rows[0] == ["状态", "源文件", "输出文件", "字数", "OCR", "图片", "耗时(ms)", "备注"]
    assert rows[1][1] == "a.md"


def test_collect_files_dir_ignores_output():
    import tempfile

    with tempfile.TemporaryDirectory() as d:
        root = Path(d)
        (root / "x.pdf").write_bytes(b"%PDF-fake")
        (root / "output").mkdir()
        (root / "output" / "old.md").write_text("skip me")
        (root / "notes.txt").write_text("hi")
        files = collect_files([root])
        names = {f.name for f in files}
        assert names == {"x.pdf", "notes.txt"}


def test_error_result_on_bad_file(tmp_path: Path):
    src = tmp_path / "broken.docx"
    src.write_bytes(b"this is not a real docx file" * 20)
    r = convert_file(src, ConvertOptions(output_dir=tmp_path / "output"))
    # 损坏文件应转为 error/warning 而不是抛异常
    assert r.status in ("error", "warning")
    assert r.message


def test_ocr_graceful_when_unavailable(tmp_path: Path, monkeypatch):
    """无 OCR 引擎时：转换不崩，图片降级为 warning"""
    import doc2md_tool.service as svc

    monkeypatch.setattr(svc, "_run_ocr_for_file", lambda src: None)
    # 造一个无文字层的图片：纯色 PNG
    from PIL import Image
    src = tmp_path / "blank.png"
    Image.new("RGB", (100, 100), "white").save(src)
    r = svc.convert_file(src, svc.ConvertOptions(output_dir=tmp_path / "out", enable_ocr=True))
    assert r.status in ("warning", "error")
    assert r.message


def test_docx_embedded_image_extracted(tmp_path: Path):
    """docx 内嵌图片：base64 不被写入 md，图片落盘，OCR 内联（若有引擎）"""
    from docx import Document
    from docx.shared import Inches
    img_p = tmp_path / "chart.png"
    _make_test_image(img_p, "ONE")

    src = tmp_path / "含图.docx"
    doc = Document()
    doc.add_heading("汇报", 0)
    doc.add_paragraph("数据见下图：")
    doc.add_picture(str(img_p), width=Inches(4))
    doc.save(src)

    out = tmp_path / "out"
    r = convert_file(src, ConvertOptions(output_dir=out))
    assert r.status == "ok", r.message
    assert r.img_count >= 1
    md = (out / "含图.md").read_text(encoding="utf-8")
    assert "base64" not in md, "base64 残留在 md 中"
    assert "![图片1](" in md
    from doc2md_tool import ocr as _ocr
    if _ocr.ocr_available():
        assert "OCR 内容" in md
    imgs = list((out / "含图_img").glob("*"))
    assert len(imgs) >= 1


def test_pptx_embedded_image_extracted(tmp_path: Path):
    """pptx 内嵌图片：从 slide 的 rels 抽出，插到对应页，OCR 内联（若有引擎）"""
    from pptx import Presentation
    from pptx.util import Inches

    img_p = tmp_path / "pic.png"
    _make_test_image(img_p, "TWO")

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "封面页"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "数据页"
    s2.shapes.add_picture(str(img_p), Inches(1), Inches(1), width=Inches(4))
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "结束页"
    src = tmp_path / "演示.pptx"
    prs.save(src)

    out = tmp_path / "out"
    r = convert_file(src, ConvertOptions(output_dir=out))
    assert r.status == "ok", r.message
    assert r.img_count >= 1
    content = (out / "演示.md").read_text(encoding="utf-8")
    assert "![图片1](演示_img/" in content
    assert (out / "演示_img").is_dir() and any((out / "演示_img").glob("s*_*"))
    from doc2md_tool import ocr as _ocr
    if _ocr.ocr_available():
        assert "OCR 内容" in content
    assert "base64" not in content


def test_docx_small_icon_dropped(tmp_path: Path):
    """小于 4KB 的装饰图/图标不进入 md"""
    from docx import Document
    from docx.shared import Inches

    from PIL import Image
    Image.new("RGB", (20, 20), "red").save(str(tmp_path / "dot.png"))
    src = tmp_path / "点.docx"
    doc = Document()
    doc.add_heading("A", 0)
    for _ in range(5):
        doc.add_paragraph("正文文字内容，足够长度避免短文本警告分支。")
    doc.add_picture(str(tmp_path / "dot.png"), width=Inches(0.2))
    doc.save(src)
    out = tmp_path / "out"
    r = convert_file(src, ConvertOptions(output_dir=out))
    md = (out / "点.md").read_text(encoding="utf-8")
    assert "base64" not in md
    assert "![图片" not in md
    assert not (out / "点_img").exists() or not (out / "点_img").iterdir()
